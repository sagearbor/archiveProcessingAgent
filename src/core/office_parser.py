from __future__ import annotations

from dataclasses import dataclass
from pathlib import Path
from typing import Dict, List, Optional

from tempfile import TemporaryDirectory

from docx import Document
from openpyxl import load_workbook
from pptx import Presentation


@dataclass
class OfficeMetadata:
    author: Optional[str] = None
    title: Optional[str] = None
    created: Optional[str] = None


class OfficeParser:
    """Parse Microsoft Office documents for text and metadata."""

    def parse_docx(self, file_path: Path) -> Dict[str, object]:
        """Parse a Word document and return structured information."""
        try:
            doc = Document(file_path)
        except Exception as exc:  # pragma: no cover - depends on external file
            raise ValueError(f"Failed to open docx: {exc}") from exc

        paragraphs = []
        headings: List[str] = []
        for p in doc.paragraphs:
            if not p.text:
                continue
            para_info = {
                "text": p.text,
                "style": p.style.name,
                "bold": any(run.bold for run in p.runs if run.text),
                "italic": any(run.italic for run in p.runs if run.text),
            }
            paragraphs.append(para_info)
            if p.style.name.startswith("Heading"):
                headings.append(p.text)

        tables = []
        for table in doc.tables:
            rows = []
            for row in table.rows:
                rows.append([cell.text for cell in row.cells])
            tables.append(rows)

        with TemporaryDirectory() as tmpdir:
            images = [str(p) for p in self.extract_images(file_path, Path(tmpdir))]

        metadata = self.get_document_metadata(file_path)
        return {
            "paragraphs": paragraphs,
            "headings": headings,
            "tables": tables,
            "images": images,
            "metadata": metadata.__dict__,
        }

    def parse_xlsx(self, file_path: Path) -> Dict[str, object]:
        """Parse an Excel workbook and return structured information."""
        try:
            wb = load_workbook(file_path, data_only=True)
        except Exception as exc:  # pragma: no cover - external
            raise ValueError(f"Failed to open xlsx: {exc}") from exc

        sheets: Dict[str, Dict[str, object]] = {}
        for sheet in wb.worksheets:
            data_rows: List[List[object]] = []
            formula_rows: List[List[Optional[str]]] = []
            comments: List[Dict[str, str]] = []
            for row in sheet.iter_rows(values_only=False):
                data_rows.append([cell.value for cell in row])
                formula_rows.append([cell.value if isinstance(cell.value, str) and cell.value.startswith("=") else None for cell in row])
                for cell in row:
                    if cell.comment:
                        comments.append({"cell": cell.coordinate, "text": cell.comment.text})
            sheets[sheet.title] = {
                "data": data_rows,
                "formulas": formula_rows,
                "comments": comments,
            }

        named_ranges = list(wb.defined_names.keys())
        metadata = self.get_document_metadata(file_path)
        return {
            "sheets": sheets,
            "named_ranges": named_ranges,
            "metadata": metadata.__dict__,
        }

    def parse_pptx(self, file_path: Path) -> Dict[str, object]:
        """Parse a PowerPoint presentation."""
        try:
            pres = Presentation(file_path)
        except Exception as exc:  # pragma: no cover - external
            raise ValueError(f"Failed to open pptx: {exc}") from exc

        slides = []
        for slide in pres.slides:
            slide_info: Dict[str, object] = {
                "layout": getattr(slide.slide_layout, "name", "Unknown"),
                "texts": [],
                "notes": slide.notes_slide.notes_text_frame.text if slide.has_notes_slide else "",
            }
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    slide_info["texts"].append(shape.text)
                if hasattr(shape, "table"):
                    table_data = []
                    for row in shape.table.rows:
                        table_data.append([cell.text for cell in row.cells])
                    slide_info.setdefault("tables", []).append(table_data)
            slides.append(slide_info)

        with TemporaryDirectory() as tmpdir:
            images = [str(p) for p in self.extract_images(file_path, Path(tmpdir))]

        metadata = self.get_document_metadata(file_path)
        return {"slides": slides, "images": images, "metadata": metadata.__dict__}

    def extract_images(self, file_path: Path, output_dir: Path) -> List[Path]:
        """Extract embedded images from Office documents."""
        output_dir.mkdir(parents=True, exist_ok=True)
        ext = file_path.suffix.lower()
        images: List[Path] = []

        if ext == ".pptx":
            pres = Presentation(file_path)
            for slide in pres.slides:
                for shape in slide.shapes:
                    if getattr(shape, "shape_type", None) == 13:  # picture
                        image = shape.image
                        fname = output_dir / image.filename
                        with open(fname, "wb") as f:
                            f.write(image.blob)
                        images.append(fname)

        elif ext == ".docx":
            doc = Document(file_path)
            rels = doc.part._rels
            for rel in rels.values():
                if "image" in rel.reltype:
                    fname = output_dir / Path(rel.target_ref).name
                    with open(fname, "wb") as f:
                        f.write(rel.target_part.blob)
                    images.append(fname)

        elif ext == ".xlsx":
            wb = load_workbook(file_path)
            for sheet in wb.worksheets:
                for img in getattr(sheet, "_images", []):
                    fname = output_dir / Path(img.path).name
                    with open(fname, "wb") as f:
                        f.write(img._data())
                    images.append(fname)

        return images

    def get_document_metadata_docx(self, doc: Document) -> OfficeMetadata:
        props = doc.core_properties
        return OfficeMetadata(
            author=props.author,
            title=props.title,
            created=str(props.created) if props.created else None,
        )

    def get_document_metadata(self, file_path: Path) -> OfficeMetadata:
        """Extract common metadata from Office documents."""
        ext = file_path.suffix.lower()
        if ext == ".docx":
            doc = Document(file_path)
            return self.get_document_metadata_docx(doc)
        if ext == ".xlsx":
            wb = load_workbook(file_path, read_only=True)
            props = wb.properties
            return OfficeMetadata(
                author=props.creator,
                title=props.title,
                created=str(props.created) if props.created else None,
            )
        if ext == ".pptx":
            pres = Presentation(file_path)
            props = pres.core_properties
            return OfficeMetadata(
                author=props.author,
                title=props.title,
                created=str(props.created) if props.created else None,
            )
        return OfficeMetadata()
