import sys
from pathlib import Path

# Ensure src package is importable
ROOT = Path(__file__).resolve().parents[1]
SRC = ROOT / "src"
if str(SRC) not in sys.path:
    sys.path.insert(0, str(SRC))
import zipfile
import tarfile
from typing import Dict
import io
from docx import Document
from openpyxl import Workbook
from pptx import Presentation
import pytest

MOCK_DIR = Path(__file__).resolve().parents[1] / "mock_data"


def _create_docx(path: Path) -> None:
    doc = Document()
    doc.add_heading("Test Document", level=1)
    doc.add_paragraph("This is a test paragraph.")
    doc.core_properties.author = "tester"
    doc.core_properties.title = "mock"
    doc.save(path)


def _create_xlsx(path: Path) -> None:
    wb = Workbook()
    ws = wb.active
    ws.title = "Sheet1"
    ws["A1"] = "data"
    wb.create_sheet("Sheet2")
    wb.save(path)


def _create_pptx(path: Path) -> None:
    prs = Presentation()
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Title"
    prs.save(path)


def _create_zip(path: Path, files: Dict[str, str]) -> None:
    with zipfile.ZipFile(path, "w") as z:
        for name, content in files.items():
            z.writestr(name, content)


def _create_tar_gz(path: Path, files: Dict[str, str]) -> None:
    with tarfile.open(path, "w:gz") as t:
        for name, content in files.items():
            tarinfo = tarfile.TarInfo(name)
            data = content.encode()
            tarinfo.size = len(data)
            t.addfile(tarinfo, io.BytesIO(data))


@pytest.fixture(scope="session", autouse=True)
def create_mock_files() -> None:
    MOCK_DIR.mkdir(exist_ok=True)
    if not (MOCK_DIR / "mock_word.docx").exists():
        _create_docx(MOCK_DIR / "mock_word.docx")
    if not (MOCK_DIR / "mock_excel.xlsx").exists():
        _create_xlsx(MOCK_DIR / "mock_excel.xlsx")
    if not (MOCK_DIR / "mock_powerpoint.pptx").exists():
        _create_pptx(MOCK_DIR / "mock_powerpoint.pptx")
    if not (MOCK_DIR / "mock_powerbi.pbix").exists():
        _create_zip(MOCK_DIR / "mock_powerbi.pbix", {"DataModel/model.bim": "{}"})
    if not (MOCK_DIR / "mock_tableau.twbx").exists():
        _create_zip(MOCK_DIR / "mock_tableau.twbx", {"workbook.xml": "<xml/>"})
    if not (MOCK_DIR / "mock_synapse.zip").exists():
        _create_zip(MOCK_DIR / "mock_synapse.zip", {"README.md": "synapse"})
    if not (MOCK_DIR / "mock_archive.zip").exists():
        _create_zip(MOCK_DIR / "mock_archive.zip", {"file.txt": "hello"})
    if not (MOCK_DIR / "mock_source.tar.gz").exists():
        _create_tar_gz(MOCK_DIR / "mock_source.tar.gz", {"src/main.py": "print()"})
